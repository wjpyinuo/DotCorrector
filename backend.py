"""
DotCorrector Backend
三级纠错流水线 + 文件解析 + 导出
"""

import os
import json
import re
from pathlib import Path
from typing import List, Dict

from PySide6.QtCore import QObject, Signal, Slot, Property, QThread


# ============================================================
#  错别字词典（第一遍）
# ============================================================

# 常见同音字 / 形近字对照  (wrong → right)
TYPO_DICT: Dict[str, str] = {
    "在坐": "在座",
    "在次": "再次",
    "在说": "再说",
    "在三": "再三",
    "不在是": "不再是",
    "不在像": "不再像",
    "不在有": "不再有",
    "以经": "已经",
    "以有": "已有",
    "以是": "已是",
    "己经": "已经",
    "自已": "自己",
    "做案": "作案",
    "做者": "作者",
    "做品": "作品",
    "做用": "作用",
    "象这样": "像这样",
    "象这种": "像这种",
    "象这种的": "像这种的",
    "好象": "好像",
    "就象": "就像",
    "真象": "真像",
    "人象": "人像",
    "画象": "画像",
    "雕象": "雕像",
    "偶象": "偶像",
    "印像": "印象",
    "想向": "想象",
    "成份": "成分",
    "分额": "份额",
    "即然": "既然",
    "既使": "即使",
    "既便": "即便",
    "记的": "记得",
    "觉的": "觉得",
    "认的": "认得",
    "懂的": "懂得",
    "晓的": "晓得",
    "值的": "值得",
    "显的": "显得",
    "听的": "听得",
    "看的见": "看得见",
    "听的到": "听得到",
    "变的": "变得",
    "说的": "说得",
    "地确": "的确",
    "必竟": "毕竟",
    "竟争": "竞争",
    "竞然": "竟然",
    "决对": "绝对",
    "因该": "应该",
    "应响": "影响",
    "影起": "引起",
    "反咉": "反映",
    "及于": "基于",
    "基与": "基于",
    "至使": "致使",
    "致今": "至今",
    "至力": "致力",
    "发辉": "发挥",
    "发杨": "发扬",
    "发奋": "发愤",
    "幅射": "辐射",
    "复盖": "覆盖",
    "破斧沉舟": "破釜沉舟",
    "一诺千斤": "一诺千金",
    "一股作气": "一鼓作气",
    "悬梁刺骨": "悬梁刺股",
    "默守成规": "墨守成规",
    "穿流不息": "川流不息",
    "迫不急待": "迫不及待",
    "谈笑风声": "谈笑风生",
    "人才倍出": "人才辈出",
    "直接了当": "直截了当",
    "走头无路": "走投无路",
    "出人投地": "出人头地",
    "因地治宜": "因地制宜",
    "淋漓尽至": "淋漓尽致",
    "按步就班": "按部就班",
    "责无旁代": "责无旁贷",
    "如法泡制": "如法炮制",
    "变本加利": "变本加厉",
    "不可思义": "不可思议",
    "不加思索": "不假思索",
    "不落巢臼": "不落窠臼",
    "不径而走": "不胫而走",
    "不醒人事": "不省人事",
    "苍海桑田": "沧海桑田",
    "出类拔粹": "出类拔萃",
    "大相径廷": "大相径庭",
    "当物之急": "当务之急",
    "独挡一面": "独当一面",
    "分道扬镖": "分道扬镳",
    "甘败下风": "甘拜下风",
    "各行其事": "各行其是",
    "功亏一匮": "功亏一篑",
    "骨梗在喉": "骨鲠在喉",
    "关怀倍至": "关怀备至",
    "鬼计多端": "诡计多端",
    "汗流夹背": "汗流浃背",
    "好高鹜远": "好高骛远",
    "和霭可亲": "和蔼可亲",
    "虎视耽耽": "虎视眈眈",
    "慌慌不安": "惶惶不安",
    "黄梁一梦": "黄粱一梦",
    "浑浑恶恶": "浑浑噩噩",
    "饥肠漉漉": "饥肠辘辘",
    "娇揉造作": "矫揉造作",
    "戒骄戒燥": "戒骄戒躁",
    "金榜提名": "金榜题名",
    "精兵减政": "精兵简政",
    "鸠占雀巢": "鸠占鹊巢",
    "开诚不公": "开诚布公",
    "开门缉盗": "开门揖盗",
    "刻划": "刻画",
    "口蜜腹箭": "口蜜腹剑",
    "苦心孤脂": "苦心孤诣",
    "烂竽充数": "滥竽充数",
    "礼上往来": "礼尚往来",
    "历兵秣马": "厉兵秣马",
    "厉精图治": "励精图治",
    "联篇累牍": "连篇累牍",
    "鳞次节比": "鳞次栉比",
    "流言非语": "流言蜚语",
    "满山遍野": "漫山遍野",
    "貌和神离": "貌合神离",
    "没精打彩": "没精打采",
    "美仑美奂": "美轮美奂",
    "迷天大罪": "弥天大罪",
    "棉里藏针": "绵里藏针",
    "明火执杖": "明火执仗",
    "名符其实": "名副其实",
    "名落深山": "名落孙山",
    "磨拳擦掌": "摩拳擦掌",
    "莫不关心": "漠不关心",
    "目不瑕接": "目不暇接",
    "弄巧成绌": "弄巧成拙",
    "旁证博引": "旁征博引",
    "披星带月": "披星戴月",
    "破锭百出": "破绽百出",
    "气极败坏": "气急败坏",
    "前扑后继": "前赴后继",
    "轻歌慢舞": "轻歌曼舞",
    "穷形尽像": "穷形尽相",
    "曲高合寡": "曲高和寡",
    "人才挤挤": "人才济济",
    "惹事生非": "惹是生非",
    "融汇贯通": "融会贯通",
    "如火如茶": "如火如荼",
    "入不付出": "入不敷出",
    "若及若离": "若即若离",
    "色厉内忍": "色厉内荏",
    "珊珊来迟": "姗姗来迟",
    "稍纵即失": "稍纵即逝",
    "身名狼藉": "声名狼藉",
    "深符众望": "深孚众望",
    "生灵图碳": "生灵涂炭",
    "拾人牙惠": "拾人牙慧",
    "食不裹腹": "食不果腹",
    "史无前列": "史无前例",
    "始作踊者": "始作俑者",
    "手不失卷": "手不释卷",
    "守株带兔": "守株待兔",
    "死心踏地": "死心塌地",
    "所向披糜": "所向披靡",
    "天翻地复": "天翻地覆",
    "挺而走险": "铤而走险",
    "同仇敌慨": "同仇敌忾",
    "投机捣把": "投机倒把",
    "图穷必见": "图穷匕见",
    "完壁归赵": "完璧归赵",
    "万古常青": "万古长青",
    "万马齐暗": "万马齐喑",
    "妄自非薄": "妄自菲薄",
    "微言精义": "微言大义",
    "尾尾动听": "娓娓动听",
    "未雨绸谬": "未雨绸缪",
    "温文而雅": "温文尔雅",
    "文过是非": "文过饰非",
    "无精打彩": "无精打采",
    "无所是从": "无所适从",
    "无忘之灾": "无妄之灾",
    "洗耳躬听": "洗耳恭听",
    "喜笑眼开": "喜笑颜开",
    "瑕不掩玉": "瑕不掩瑜",
    "相辅相承": "相辅相成",
    "向偶而泣": "向隅而泣",
    "消声匿迹": "销声匿迹",
    "心狠手棘": "心狠手辣",
    "心恢意冷": "心灰意冷",
    "心惊胆颤": "心惊胆战",
    "信口词黄": "信口雌黄",
    "兴高彩烈": "兴高采烈",
    "形消骨立": "形销骨立",
    "修茸一新": "修葺一新",
    "诩诩如生": "栩栩如生",
    "宣宾夺主": "喧宾夺主",
    "循规蹈距": "循规蹈矩",
    "循序渐近": "循序渐进",
    "言简意骇": "言简意赅",
    "淹没无闻": "湮没无闻",
    "养尊处忧": "养尊处优",
    "摇摇欲堕": "摇摇欲坠",
    "一愁莫展": "一筹莫展",
    "一视同人": "一视同仁",
    "英雄倍出": "英雄辈出",
    "应接不遐": "应接不暇",
    "拥容华贵": "雍容华贵",
    "庸人自忧": "庸人自扰",
    "忧心冲冲": "忧心忡忡",
    "有持无恐": "有恃无恐",
    "有口皆杯": "有口皆碑",
    "余音绕粱": "余音绕梁",
    "原形必露": "原形毕露",
    "缘木求渔": "缘木求鱼",
    "怨天犹人": "怨天尤人",
    "越俎代疱": "越俎代庖",
    "运筹维幄": "运筹帷幄",
    "沾轻怕重": "拈轻怕重",
    "斩钉接铁": "斩钉截铁",
    "仗义直言": "仗义执言",
    "招摇装骗": "招摇撞骗",
    "针贬时弊": "针砭时弊",
    "真知卓见": "真知灼见",
    "振聋发馈": "振聋发聩",
    "指高气扬": "趾高气扬",
    "置若网闻": "置若罔闻",
    "中流抵柱": "中流砥柱",
    "众口烁金": "众口铄金",
    "珠联壁合": "珠联璧合",
    "专心至志": "专心致志",
    "壮志难筹": "壮志难酬",
    "捉襟见胄": "捉襟见肘",
    "自抱自弃": "自暴自弃",
    "自名得意": "自鸣得意",
    "坐地分脏": "坐地分赃",
    "坐想其成": "坐享其成",
}

def dict_correct(text: str) -> tuple:
    """第一遍：词典纠错，返回 (changes, corrected_text)"""
    # 收集所有匹配
    matches = []
    for wrong, right in TYPO_DICT.items():
        start = 0
        while True:
            idx = text.find(wrong, start)
            if idx == -1:
                break
            matches.append({
                "original": wrong,
                "corrected": right,
                "position": idx,
                "pass_name": "词典",
                "pass_index": 1,
                "accepted": True,
            })
            start = idx + len(wrong)

    if not matches:
        return [], text

    # 按位置正序、长词优先排序，去重重叠
    matches.sort(key=lambda m: (m["position"], -len(m["original"])))
    seen = set()
    unique = []
    for m in matches:
        key = (m["position"], m["original"])
        if key not in seen:
            seen.add(key)
            unique.append(m)

    # 正序替换 + offset 追踪，确保位置准确
    result, corrected = [], text
    offset = 0
    for m in unique:
        adj_pos = m["position"] + offset
        orig = m["original"]
        if adj_pos + len(orig) <= len(corrected) and corrected[adj_pos:adj_pos + len(orig)] == orig:
            corrected = corrected[:adj_pos] + m["corrected"] + corrected[adj_pos + len(orig):]
            result.append({**m, "position": adj_pos})
            offset += len(m["corrected"]) - len(orig)
    return result, corrected


# ============================================================
#  pycorrector（第二遍）
# ============================================================

def pycorrector_correct(text: str) -> tuple:
    """第二遍：pycorrector 纠错"""
    try:
        from pycorrector import Corrector
    except ImportError:
        return [], text

    cr = Corrector()
    result = cr.correct(text)
    details = result.get("details", [])

    changes = []
    # pycorrector 返回的位置可能有偏移，按位置降序处理
    details_sorted = sorted(details, key=lambda d: d.get("begin_idx", 0), reverse=True)

    corrected_text = text
    for det in details_sorted:
        pos = det.get("begin_idx", 0)
        wrong = det.get("error_word", "")
        right = det.get("correct_word", "")
        if not wrong or not right:
            continue
        # 验证位置
        if pos + len(wrong) <= len(corrected_text) and corrected_text[pos:pos + len(wrong)] == wrong:
            corrected_text = corrected_text[:pos] + right + corrected_text[pos + len(wrong):]
            changes.append({
                "original": wrong,
                "corrected": right,
                "position": pos,
                "pass_name": "pycorrector",
                "pass_index": 2,
                "accepted": True,
            })

    # 因为我们是降序替换，位置已经不正确了，需要正序重新计算
    # 重新从原文算一遍
    changes_final = []
    offset = 0
    for ch in sorted(changes, key=lambda c: c["position"]):
        ch_copy = dict(ch)
        ch_copy["position"] = ch["position"] + offset
        changes_final.append(ch_copy)
        offset += len(ch["corrected"]) - len(ch["original"])

    return changes_final, corrected_text


# ============================================================
#  大模型 AI 精校（第三遍）
# ============================================================

LLM_PROMPT = """你是一个中文文本纠错专家。请仔细检查以下文本中的错别字、语法错误和用词不当。

要求：
1. 只纠错，不改写，保留原文风格和意思
2. 用 JSON 格式返回，包含 corrections 数组
3. 每个 correction 包含：position（字符位置），original（原文片段），corrected（修正后片段）
4. 如果没有错误，返回空数组

文本：
{text}

请返回 JSON 格式（不要包含 markdown 代码块标记）：
{{"corrections": [{{"position": 0, "original": "xxx", "corrected": "yyy"}}]}}
"""


def llm_correct(text: str, api_key: str, api_base: str = "https://api.mimo.ai/v1",
                provider: str = "mimo") -> tuple:
    """第三遍：大模型 API 精校"""
    if not api_key:
        return [], text

    try:
        import requests
    except ImportError:
        return [], text

    model_name = "mimo-v2-pro" if provider == "mimo" else "deepseek-chat"

    # 长文本分段（每段最多 2000 字）
    chunk_size = 2000
    chunks = []
    for i in range(0, len(text), chunk_size):
        chunks.append((i, text[i:i + chunk_size]))

    all_changes = []
    corrected_text = text

    for offset, chunk in chunks:
        try:
            resp = requests.post(
                f"{api_base}/chat/completions",
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": model_name,
                    "messages": [
                        {"role": "system", "content": "你是中文文本纠错专家。"},
                        {"role": "user", "content": LLM_PROMPT.format(text=chunk)},
                    ],
                    "temperature": 0,
                },
                timeout=30,
            )

            if resp.status_code != 200:
                continue

            data = resp.json()
            content = data["choices"][0]["message"]["content"].strip()

            # 清理可能的 markdown 代码块
            content = re.sub(r"```json\s*", "", content)
            content = re.sub(r"```\s*$", "", content)

            result = json.loads(content)
            corrections = result.get("corrections", [])

            # 按位置降序替换
            corrections_sorted = sorted(corrections, key=lambda c: c.get("position", 0), reverse=True)

            for cor in corrections_sorted:
                pos = cor.get("position", 0) + offset
                original = cor.get("original", "")
                corrected = cor.get("corrected", "")
                if not original or not corrected:
                    continue
                if pos + len(original) <= len(corrected_text) and corrected_text[pos:pos + len(original)] == original:
                    corrected_text = corrected_text[:pos] + corrected + corrected_text[pos + len(original):]
                    all_changes.append({
                        "original": original,
                        "corrected": corrected,
                        "position": pos,
                        "pass_name": "AI精校",
                        "pass_index": 3,
                        "accepted": True,
                    })

        except Exception:
            continue

    # 重新计算位置（正序）
    changes_final = []
    offset_calc = 0
    for ch in sorted(all_changes, key=lambda c: c["position"]):
        ch_copy = dict(ch)
        ch_copy["position"] = ch["position"] + offset_calc
        changes_final.append(ch_copy)
        offset_calc += len(ch["corrected"]) - len(ch["original"])

    return changes_final, corrected_text


# ============================================================
#  文件解析 / 导出
# ============================================================

def extract_text_from_file(filepath: str) -> tuple:
    """从文件提取文本，返回 (segments, file_type)
    segments: list of {"text": str, "ref": any}  ref 用于导出时定位
    """
    ext = Path(filepath).suffix.lower()

    if ext == ".txt":
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return [{"text": text, "ref": None}], "txt"

    elif ext == ".docx":
        from docx import Document
        doc = Document(filepath)
        segments = []
        for para in doc.paragraphs:
            if para.text.strip():
                segments.append({"text": para.text, "ref": para})
        return segments, "docx"

    elif ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(filepath)
        segments = []
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str) and cell.value.strip():
                        segments.append({"text": cell.value, "ref": cell})
        return segments, "xlsx"

    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(filepath)
        segments = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            segments.append({"text": para.text, "ref": (shape.text_frame, para)})
        return segments, "pptx"

    else:
        raise ValueError(f"不支持的格式: {ext}")


def export_corrected_file(segments: List[Dict], file_type: str, output_path: str, original_path: str):
    """导出修正后的文件"""
    if file_type == "txt":
        with open(output_path, "w", encoding="utf-8") as f:
            for i, seg in enumerate(segments):
                if i > 0:
                    f.write("\n\n")
                f.write(seg["corrected_text"])

    elif file_type == "docx":
        from docx import Document
        doc = Document(original_path)
        for seg in segments:
            ref = seg.get("ref")
            if ref and hasattr(ref, "text"):
                # 清除原段落内容并写入修正文本
                for run in ref.runs:
                    run.text = ""
                if ref.runs:
                    ref.runs[0].text = seg["corrected_text"]
                else:
                    ref.add_run(seg["corrected_text"])
        doc.save(output_path)

    elif file_type == "xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(original_path)
        for seg in segments:
            ref = seg.get("ref")
            if ref:
                ref.value = seg["corrected_text"]
        wb.save(output_path)

    elif file_type == "pptx":
        from pptx import Presentation
        prs = Presentation(original_path)
        seg_idx = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip() and seg_idx < len(segments):
                            for run in para.runs:
                                run.text = ""
                            if para.runs:
                                para.runs[0].text = segments[seg_idx]["corrected_text"]
                            seg_idx += 1
        prs.save(output_path)


# ============================================================
#  纠错工作线程
# ============================================================

class CorrectWorker(QThread):
    progressChanged = Signal(int)
    statusChanged = Signal(str)
    segmentDone = Signal(int, str, list)  # seg_index, corrected_text, changes
    finished_result = Signal(list)  # all segments with results

    def __init__(self, files: list, use_dict=True, use_pycorrector=True, use_llm=False,
                 api_key="", api_base="https://api.mimo.ai/v1", provider="mimo"):
        super().__init__()
        self.files = files
        self.use_dict = use_dict
        self.use_pycorrector = use_pycorrector
        self.use_llm = use_llm
        self.api_key = api_key
        self.api_base = api_base
        self.provider = provider

    def run(self):
        all_results = []

        for fi, filepath in enumerate(self.files):
            try:
                segments, file_type = extract_text_from_file(filepath)
            except Exception as e:
                self.statusChanged.emit(f"解析失败: {Path(filepath).name} - {e}")
                continue

            file_results = []
            total_segs = len(segments)

            for si, seg in enumerate(segments):
                text = seg["text"]
                all_changes = []
                current_text = text

                # 第一遍：词典
                if self.use_dict:
                    self.statusChanged.emit(f"[{fi+1}/{len(self.files)}] 词典纠错... ({si+1}/{total_segs})")
                    dict_changes, current_text = dict_correct(current_text)
                    all_changes.extend(dict_changes)

                # 第二遍：pycorrector
                if self.use_pycorrector:
                    self.statusChanged.emit(f"[{fi+1}/{len(self.files)}] pycorrector 纠错... ({si+1}/{total_segs})")
                    pyc_changes, current_text = pycorrector_correct(current_text)
                    all_changes.extend(pyc_changes)

                # 第三遍：大模型
                if self.use_llm and self.api_key:
                    self.statusChanged.emit(f"[{fi+1}/{len(self.files)}] AI 精校... ({si+1}/{total_segs})")
                    llm_changes, current_text = llm_correct(current_text, self.api_key, self.api_base, self.provider)
                    all_changes.extend(llm_changes)

                result = {
                    "file": filepath,
                    "file_type": file_type,
                    "seg_index": si,
                    "original_text": text,
                    "corrected_text": current_text,
                    "changes": all_changes,
                    "ref": seg["ref"],
                }
                file_results.append(result)

                pct = int(((fi * total_segs + si + 1) / (len(self.files) * total_segs)) * 100)
                self.progressChanged.emit(pct)
                self.segmentDone.emit(si, current_text, all_changes)

            all_results.extend(file_results)

        self.statusChanged.emit(f"完成！共处理 {len(self.files)} 个文件")
        self.finished_result.emit(all_results)


# ============================================================
#  主后端对象（暴露给 QML）
# ============================================================

class Backend(QObject):
    progressChanged = Signal(int)
    statusChanged = Signal(str)
    busyChanged = Signal(bool)
    resultsChanged = Signal()
    reportTextChanged = Signal()
    settingsChanged = Signal()

    def __init__(self):
        super().__init__()
        self._progress = 0
        self._status = "就绪"
        self._busy = False
        self._results = []  # list of result dicts
        self._report_text = ""
        self._worker = None
        self._files = []

        # 设置
        self._use_dict = True
        self._use_pycorrector = True
        self._use_llm = False
        self._api_key = ""
        self._api_provider = "mimo"
        self._api_base = "https://api.mimo.ai/v1"

        self._API_ENDPOINTS = {
            "mimo": "https://api.mimo.ai/v1",
            "deepseek": "https://api.deepseek.com/v1",
        }

    # ---- progress ----
    def _get_progress(self):
        return self._progress

    def _set_progress(self, v):
        if self._progress != v:
            self._progress = v
            self.progressChanged.emit(v)

    progress = Property(int, _get_progress, _set_progress, notify=progressChanged)

    # ---- status ----
    def _get_status(self):
        return self._status

    def _set_status(self, v):
        if self._status != v:
            self._status = v
            self.statusChanged.emit(v)

    status = Property(str, _get_status, _set_status, notify=statusChanged)

    # ---- busy ----
    def _get_busy(self):
        return self._busy

    def _set_busy(self, v):
        if self._busy != v:
            self._busy = v
            self.busyChanged.emit(v)

    busy = Property(bool, _get_busy, _set_busy, notify=busyChanged)

    # ---- reportText ----
    def _get_report_text(self):
        return self._report_text

    reportText = Property(str, _get_report_text, notify=reportTextChanged)

    # ---- results (JSON string for QML) ----
    def _get_results_json(self):
        if not self._results:
            return "[]"
        # 返回简化后的结果给 QML
        simplified = []
        for r in self._results:
            simplified.append({
                "file": Path(r["file"]).name,
                "file_path": r["file"],
                "file_type": r["file_type"],
                "seg_index": r["seg_index"],
                "original_text": r["original_text"],
                "corrected_text": r["corrected_text"],
                "changes": r["changes"],
            })
        return json.dumps(simplified, ensure_ascii=False)

    resultsJson = Property(str, _get_results_json, notify=resultsChanged)

    # ---- 总改动数 ----
    def _get_total_changes(self):
        return sum(len(r.get("changes", [])) for r in self._results)

    totalChanges = Property(int, _get_total_changes, notify=resultsChanged)

    # ---- 按遍数统计 ----
    def _get_dict_changes(self):
        return sum(1 for r in self._results for c in r.get("changes", []) if c["pass_index"] == 1)

    dictChanges = Property(int, _get_dict_changes, notify=resultsChanged)

    def _get_pyc_changes(self):
        return sum(1 for r in self._results for c in r.get("changes", []) if c["pass_index"] == 2)

    pycChanges = Property(int, _get_pyc_changes, notify=resultsChanged)

    def _get_llm_changes(self):
        return sum(1 for r in self._results for c in r.get("changes", []) if c["pass_index"] == 3)

    llmChanges = Property(int, _get_llm_changes, notify=resultsChanged)

    # ---- 设置属性 ----
    def _get_use_dict(self):
        return self._use_dict

    def _set_use_dict(self, v):
        self._use_dict = v
        self.settingsChanged.emit()

    useDict = Property(bool, _get_use_dict, _set_use_dict, notify=settingsChanged)

    def _get_use_pycorrector(self):
        return self._use_pycorrector

    def _set_use_pycorrector(self, v):
        self._use_pycorrector = v
        self.settingsChanged.emit()

    usePycorrector = Property(bool, _get_use_pycorrector, _set_use_pycorrector, notify=settingsChanged)

    def _get_use_llm(self):
        return self._use_llm

    def _set_use_llm(self, v):
        self._use_llm = v
        self.settingsChanged.emit()

    useLlm = Property(bool, _get_use_llm, _set_use_llm, notify=settingsChanged)

    def _get_api_key(self):
        return self._api_key

    def _set_api_key(self, v):
        self._api_key = v
        self.settingsChanged.emit()

    apiKey = Property(str, _get_api_key, _set_api_key, notify=settingsChanged)

    def _get_api_provider(self):
        return self._api_provider

    def _set_api_provider(self, v):
        if self._api_provider != v:
            self._api_provider = v
            self._api_base = self._API_ENDPOINTS.get(v, self._api_base)
            self.settingsChanged.emit()

    apiProvider = Property(str, _get_api_provider, _set_api_provider, notify=settingsChanged)

    # ---- QML 可调用方法 ----

    @Slot(list)
    def startCorrect(self, files):
        if self._busy:
            return

        # 至少启用一种纠错
        if not self._use_dict and not self._use_pycorrector and not self._use_llm:
            self._set_status("请至少启用一种纠错引擎")
            return

        self._files = []
        for f in files:
            path = f.toString().replace("file:///", "") if isinstance(f, str) else str(f)
            if os.path.isfile(path):
                self._files.append(path)

        if not self._files:
            self._set_status("未找到有效文件")
            return

        self._set_busy(True)
        self._set_progress(0)
        self._set_status(f"正在处理 {len(self._files)} 个文件...")
        self._results = []

        self._worker = CorrectWorker(
            self._files,
            use_dict=self._use_dict,
            use_pycorrector=self._use_pycorrector,
            use_llm=self._use_llm,
            api_key=self._api_key,
            api_base=self._api_base,
            provider=self._api_provider,
        )
        self._worker.progressChanged.connect(self._set_progress)
        self._worker.statusChanged.connect(self._set_status)
        self._worker.finished_result.connect(self._on_done)
        self._worker.start()

    def _on_done(self, results):
        self._results = results
        self._set_busy(False)
        self._set_progress(100)
        self.resultsChanged.emit()
        self._generate_report()
        self._set_status(f"完成！共 {self.totalChanges} 处修改")

    def _generate_report(self):
        lines = ["=== 墨正 DotCorrector 纠错报告 ===", ""]
        lines.append(f"文件数: {len(set(r['file'] for r in self._results))}")
        lines.append(f"总修改: {self.totalChanges} 处")
        lines.append(f"  词典纠错: {self.dictChanges} 处")
        lines.append(f"  pycorrector: {self.pycChanges} 处")
        lines.append(f"  AI 精校: {self.llmChanges} 处")
        lines.append("")

        for r in self._results:
            fname = Path(r["file"]).name
            lines.append(f"--- {fname} (段落 {r['seg_index'] + 1}) ---")
            lines.append(f"原文: {r['original_text'][:100]}...")
            if r["changes"]:
                for ch in r["changes"]:
                    lines.append(f"  [{ch['pass_name']}] {ch['original']} → {ch['corrected']}")
            else:
                lines.append("  无修改")
            lines.append("")

        self._report_text = "\n".join(lines)
        self.reportTextChanged.emit()

    @Slot(int, int, bool)
    def setChangeAccepted(self, seg_index, change_index, accepted):
        """接受/拒绝某个改动"""
        if 0 <= seg_index < len(self._results):
            changes = self._results[seg_index].get("changes", [])
            if 0 <= change_index < len(changes):
                changes[change_index]["accepted"] = accepted
                self.resultsChanged.emit()

    @Slot(int)
    def acceptAllChanges(self, pass_index=0):
        """接受所有改动（0=全部）"""
        for r in self._results:
            for ch in r.get("changes", []):
                if pass_index == 0 or ch["pass_index"] == pass_index:
                    ch["accepted"] = True
        self.resultsChanged.emit()

    @Slot(int)
    def rejectAllChanges(self, pass_index=0):
        """拒绝所有改动（0=全部）"""
        for r in self._results:
            for ch in r.get("changes", []):
                if pass_index == 0 or ch["pass_index"] == pass_index:
                    ch["accepted"] = False
        self.resultsChanged.emit()

    @Slot(str)
    def exportResults(self, output_dir):
        """导出修正后的文件"""
        if not self._results:
            self._set_status("没有可导出的结果")
            return

        output_dir = output_dir.toString().replace("file:///", "") if isinstance(output_dir, str) else output_dir
        os.makedirs(output_dir, exist_ok=True)

        # 按文件分组
        files_groups = {}
        for r in self._results:
            fp = r["file"]
            if fp not in files_groups:
                files_groups[fp] = {"type": r["file_type"], "segments": []}
            # 只应用被接受的改动
            accepted_text = r["original_text"]
            for ch in sorted(r["changes"], key=lambda c: c["position"], reverse=True):
                if ch["accepted"]:
                    pos = ch["position"]
                    if pos + len(ch["original"]) <= len(accepted_text):
                        accepted_text = accepted_text[:pos] + ch["corrected"] + accepted_text[pos + len(ch["original"]):]
            files_groups[fp]["segments"].append({
                "corrected_text": accepted_text,
                "ref": r.get("ref"),
            })

        exported = 0
        for filepath, group in files_groups.items():
            stem = Path(filepath).stem
            ext = Path(filepath).suffix
            output_path = os.path.join(output_dir, f"{stem}_corrected{ext}")
            try:
                export_corrected_file(group["segments"], group["type"], output_path, filepath)
                exported += 1
            except Exception as e:
                self._set_status(f"导出失败: {e}")
                return

        self._set_status(f"已导出 {exported} 个文件到 {output_dir}")

    @Slot(str, result=str)
    def readFileContent(self, filepath):
        """读取文件内容（预览用）"""
        filepath = filepath.toString().replace("file:///", "") if isinstance(filepath, str) else filepath
        try:
            segments, _ = extract_text_from_file(filepath)
            return "\n\n".join(s["text"] for s in segments)
        except Exception as e:
            return f"读取失败: {e}"
